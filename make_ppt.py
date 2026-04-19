"""
Generate HRGA Activity App PowerPoint Presentation
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import os

# ── Brand Colours ──────────────────────────────────────────────
BLUE_DEEP   = RGBColor(0x07, 0x3b, 0x73)
BLUE_MAIN   = RGBColor(0x0f, 0x5f, 0xb8)
BLUE_SOFT   = RGBColor(0xe8, 0xf2, 0xff)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
INK         = RGBColor(0x1f, 0x29, 0x37)
MUTED       = RGBColor(0x64, 0x74, 0x8b)
AMBER       = RGBColor(0xf5, 0x9e, 0x0b)
GREEN       = RGBColor(0x16, 0xa3, 0x4a)
RED         = RGBColor(0xdc, 0x26, 0x26)
PURPLE      = RGBColor(0x7c, 0x3a, 0xed)
LIGHT_GRAY  = RGBColor(0xf5, 0xf8, 0xfc)
BORDER      = RGBColor(0xe5, 0xed, 0xf7)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]   # blank


# ── helpers ────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width or 1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, l, t, w, h,
                size=12, bold=False, color=INK,
                align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_para(tf, text, size=11, bold=False, color=INK,
             align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def gradient_bg(slide, top_color=BLUE_DEEP, bot_color=BLUE_MAIN):
    """Full-slide gradient rectangle (simulated with two rects)."""
    add_rect(slide, 0, 0, 13.33, 3.75, fill=top_color)
    add_rect(slide, 0, 3.75, 13.33, 3.75, fill=bot_color)


def section_header(slide, title, subtitle=""):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=BLUE_SOFT)
    add_rect(slide, 0, 0, 13.33, 1.2, fill=BLUE_MAIN)
    add_rect(slide, 0, 1.2, 0.06, 6.3, fill=BLUE_DEEP)
    add_textbox(slide, title, 0.2, 0.28, 12, 0.7,
                size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle, 0.2, 0.88, 12, 0.4,
                    size=13, bold=False, color=RGBColor(0xc7, 0xdc, 0xf8),
                    align=PP_ALIGN.LEFT)


def feature_card(slide, icon, title, bullets, l, t, w=5.8, h=2.8,
                 accent=BLUE_MAIN):
    add_rect(slide, l, t, w, h, fill=WHITE, line_color=BORDER, line_width=0.75)
    add_rect(slide, l, t, 0.045, h, fill=accent)
    add_textbox(slide, icon + "  " + title,
                l+0.12, t+0.12, w-0.2, 0.4,
                size=13, bold=True, color=BLUE_DEEP)
    txb = slide.shapes.add_textbox(
        Inches(l+0.12), Inches(t+0.58), Inches(w-0.22), Inches(h-0.7))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = "•  " + b
        run.font.size = Pt(10.5)
        run.font.color.rgb = INK


def screenshot_placeholder(slide, l, t, w, h, label="Screenshot"):
    add_rect(slide, l, t, w, h, fill=RGBColor(0xee, 0xf4, 0xff),
             line_color=BLUE_MAIN, line_width=1)
    # inner dashed border effect
    add_rect(slide, l+0.05, t+0.05, w-0.1, h-0.1,
             line_color=RGBColor(0xb8, 0xd4, 0xf5), line_width=0.5)
    add_textbox(slide, "📷", l, t + h/2 - 0.35, w, 0.4,
                size=24, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, l, t + h/2 + 0.05, w, 0.35,
                size=10, color=MUTED, align=PP_ALIGN.CENTER)


def stat_box(slide, value, label, l, t, w=2.5, h=1.3, accent=BLUE_MAIN):
    add_rect(slide, l, t, w, h, fill=WHITE, line_color=BORDER, line_width=0.75)
    add_rect(slide, l, t, w, 0.06, fill=accent)
    add_textbox(slide, value, l, t+0.15, w, 0.55,
                size=28, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, l, t+0.68, w, 0.45,
                size=10, color=MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 – COVER
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
gradient_bg(slide)
add_rect(slide, 0, 0, 13.33, 7.5, fill=None)

# Decorative circles
add_rect(slide, 9.5, -1.2, 5, 5, fill=RGBColor(0x0a, 0x48, 0x8c))
add_rect(slide, 10.8, 4.5, 4, 4, fill=RGBColor(0x06, 0x30, 0x5e))
add_rect(slide, -1, 5, 3.5, 3.5, fill=RGBColor(0x06, 0x30, 0x5e))

# Logo area
add_rect(slide, 1.0, 1.6, 1.2, 0.9, fill=WHITE)
add_textbox(slide, "SCM", 1.02, 1.65, 1.15, 0.8,
            size=22, bold=True, color=BLUE_MAIN, align=PP_ALIGN.CENTER)

# Title block
add_textbox(slide, "HRGA Activity", 2.5, 1.5, 9, 0.85,
            size=46, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_textbox(slide, "Sistem Manajemen Aktivitas Harian & Proyek",
            2.5, 2.42, 9, 0.55,
            size=18, bold=False, color=RGBColor(0xc7, 0xdc, 0xf8),
            align=PP_ALIGN.LEFT)

# Divider
add_rect(slide, 1.0, 3.2, 11.33, 0.04, fill=RGBColor(0x4d, 0x96, 0xe0))

# Subtitle row
add_textbox(slide, "PT Sulawesi Cahaya Mineral  |  Departemen HRGA",
            1.0, 3.35, 10, 0.45,
            size=14, color=RGBColor(0xb8, 0xd4, 0xf8),
            align=PP_ALIGN.LEFT)

# Feature pills
pills = ["Daily Activity", "Kalender", "Project Agile", "Dashboard", "User Management"]
px = 1.0
for p in pills:
    pw = len(p) * 0.13 + 0.5
    add_rect(slide, px, 4.1, pw, 0.42, fill=RGBColor(0x14, 0x5a, 0xa8),
             line_color=RGBColor(0x4d, 0x96, 0xe0), line_width=0.5)
    add_textbox(slide, p, px+0.05, 4.13, pw-0.05, 0.35,
                size=11, color=WHITE, align=PP_ALIGN.CENTER)
    px += pw + 0.22

# Bottom bar
add_rect(slide, 0, 6.8, 13.33, 0.7, fill=RGBColor(0x05, 0x28, 0x4d))
add_textbox(slide, "Presentasi Aplikasi  •  HRGA Activity v1.0  •  2024",
            0, 6.85, 13.33, 0.5,
            size=11, color=RGBColor(0x8b, 0xb4, 0xd8), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 – AGENDA
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
add_rect(slide, 0, 0, 13.33, 1.35, fill=BLUE_DEEP)
add_rect(slide, 0, 1.35, 0.06, 6.15, fill=AMBER)
add_textbox(slide, "Daftar Isi", 0.3, 0.32, 12, 0.7,
            size=30, bold=True, color=WHITE)
add_textbox(slide, "Overview presentasi fitur-fitur aplikasi",
            0.3, 0.92, 12, 0.4, size=13, color=RGBColor(0xc7, 0xdc, 0xf8))

items = [
    ("01", "Latar Belakang & Tujuan Aplikasi"),
    ("02", "Arsitektur & Teknologi"),
    ("03", "Fitur Login & Autentikasi"),
    ("04", "Dashboard & Statistik"),
    ("05", "Daily Activity Management"),
    ("06", "Kalender & Event"),
    ("07", "Project Agile (Sprint & Kanban)"),
    ("08", "Manajemen User (Admin)"),
    ("09", "Ringkasan & Manfaat"),
]

cols = [items[:5], items[5:]]
cx = [0.4, 6.9]
for ci, col in enumerate(cols):
    for ri, (num, label) in enumerate(col):
        ty = 1.6 + ri * 1.0
        tx = cx[ci]
        add_rect(slide, tx, ty, 0.58, 0.62, fill=BLUE_MAIN)
        add_textbox(slide, num, tx, ty+0.1, 0.58, 0.45,
                    size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, tx+0.65, ty+0.08, 5.5, 0.5, fill=WHITE,
                 line_color=BORDER, line_width=0.5)
        add_textbox(slide, label, tx+0.75, ty+0.12, 5.3, 0.38,
                    size=12, bold=False, color=INK)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 – LATAR BELAKANG
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
section_header(slide, "01  Latar Belakang & Tujuan",
               "Mengapa HRGA Activity dibangun?")
add_rect(slide, 0.3, 1.4, 8.1, 5.75, fill=WHITE, line_color=BORDER, line_width=0.75)
add_rect(slide, 8.6, 1.4, 4.5, 5.75, fill=WHITE, line_color=BORDER, line_width=0.75)

txb = slide.shapes.add_textbox(Inches(0.45), Inches(1.58), Inches(7.9), Inches(5.4))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
add_para(tf, "Permasalahan yang Diselesaikan", size=14, bold=True, color=BLUE_DEEP)
problems = [
    "Pencatatan aktivitas karyawan masih manual di spreadsheet",
    "Tidak ada visibilitas real-time progres pekerjaan tim",
    "Rapat dan event kalender tersebar di berbagai media",
    "Manajemen proyek tidak terintegrasi dengan aktivitas harian",
    "Pelaporan performa karyawan membutuhkan waktu lama",
    "Tidak ada sistem prioritas untuk pekerjaan harian",
]
for pr in problems:
    add_para(tf, "✗  " + pr, size=11, color=RED, space_before=5)

add_para(tf, "", size=6)
add_para(tf, "Solusi yang Ditawarkan", size=14, bold=True, color=GREEN)
solutions = [
    "Platform terpadu untuk tracking aktivitas harian",
    "Dashboard real-time dengan chart dan statistik",
    "Kalender terintegrasi untuk event & pengingat",
    "Modul Agile (Sprint & Kanban) untuk manajemen proyek",
    "Laporan otomatis progres karyawan untuk admin",
]
for s in solutions:
    add_para(tf, "✓  " + s, size=11, color=GREEN, space_before=5)

txb2 = slide.shapes.add_textbox(Inches(8.75), Inches(1.58), Inches(4.2), Inches(5.4))
txb2.word_wrap = True
tf2 = txb2.text_frame
tf2.word_wrap = True
add_para(tf2, "Target Pengguna", size=14, bold=True, color=BLUE_DEEP)
add_para(tf2, "", size=4)

users_info = [
    ("👤 Karyawan", "Input & pantau aktivitas harian, progres kerja, & kalender"),
    ("🔑 Admin / HRGA", "Kelola semua karyawan, laporan performa, manajemen user"),
    ("📋 Project Lead", "Buat proyek, sprint, dan assign task ke anggota tim"),
]
for icon_label, desc in users_info:
    add_para(tf2, icon_label, size=12, bold=True, color=BLUE_MAIN, space_before=8)
    add_para(tf2, desc, size=10, color=MUTED, space_before=2)

add_para(tf2, "", size=8)
add_para(tf2, "Stack Teknologi", size=14, bold=True, color=BLUE_DEEP)
add_para(tf2, "", size=4)
techs = ["Laravel 12 (PHP)", "AdminLTE 3 (UI Framework)", "MySQL Database",
         "Bootstrap 4", "Chart.js", "FullCalendar.js"]
for t in techs:
    add_para(tf2, "▸  " + t, size=10.5, color=INK, space_before=3)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 – ARSITEKTUR
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
section_header(slide, "02  Arsitektur & Teknologi",
               "Struktur MVC dan komponen sistem")
add_rect(slide, 0.3, 1.35, 12.73, 5.85, fill=WHITE, line_color=BORDER, line_width=0.75)

layers = [
    ("🌐 Layer Presentasi", BLUE_MAIN,
     ["Blade Templates (Laravel)", "AdminLTE 3.2 + Bootstrap 4",
      "Chart.js (statistik visual)", "FullCalendar (kalender interaktif)"]),
    ("⚙️  Layer Bisnis", RGBColor(0x7c, 0x3a, 0xed),
     ["9 Controllers MVC", "Role-based Access Control",
      "Form Validation & Policy", "Service Layer (Activity comments)"]),
    ("🗄️  Layer Data", GREEN,
     ["6 Eloquent Models", "MySQL via Laravel Eloquent",
      "Database Migrations", "Eloquent Relationships"]),
    ("🔐 Keamanan", AMBER,
     ["Laravel Auth Middleware", "CSRF Protection",
      "Session Management", "Role Guard (Admin/User)"]),
]

lx = 0.55
for title, accent, points in layers:
    add_rect(slide, lx, 1.55, 2.95, 5.45, fill=BLUE_SOFT,
             line_color=accent, line_width=1)
    add_rect(slide, lx, 1.55, 2.95, 0.55, fill=accent)
    add_textbox(slide, title, lx+0.1, 1.58, 2.75, 0.48,
                size=11, bold=True, color=WHITE)
    txb = slide.shapes.add_textbox(
        Inches(lx+0.12), Inches(2.22), Inches(2.72), Inches(4.6))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for pt in points:
        p = tf.paragraphs[0] if pt == points[0] else tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = "• " + pt
        run.font.size = Pt(10.5)
        run.font.color.rgb = INK
    lx += 3.12

# Arrow connectors
for ax in [3.42, 6.54, 9.66]:
    add_textbox(slide, "→", ax, 3.55, 0.28, 0.45,
                size=18, bold=True, color=MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 – LOGIN
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
section_header(slide, "03  Login & Autentikasi",
               "Sistem masuk dengan role-based access control")

# Left – screenshot placeholder
screenshot_placeholder(slide, 0.3, 1.4, 5.8, 5.8, "Tampilan Halaman Login")

# Right – info
add_rect(slide, 6.4, 1.4, 6.65, 5.8, fill=WHITE, line_color=BORDER, line_width=0.75)
txb = slide.shapes.add_textbox(Inches(6.6), Inches(1.6), Inches(6.3), Inches(5.4))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
add_para(tf, "Fitur Autentikasi", size=16, bold=True, color=BLUE_DEEP)
add_para(tf, "", size=4)
feats = [
    ("Email + Password Login", "Autentikasi aman menggunakan kredensial unik"),
    ("Role-Based Access", "Admin memiliki akses penuh; Karyawan hanya data sendiri"),
    ("Session Management", "Session diregenerasi setelah login untuk keamanan"),
    ("CSRF Protection", "Semua form dilindungi token CSRF Laravel"),
    ("Ganti Password", "Setiap user dapat mengubah password kapan saja"),
    ("Reset Password (Admin)", "Admin dapat mereset password karyawan"),
]
for label, desc in feats:
    add_para(tf, "▶  " + label, size=12, bold=True, color=BLUE_MAIN, space_before=8)
    add_para(tf, "    " + desc, size=10.5, color=MUTED, space_before=1)

add_para(tf, "", size=6)
add_para(tf, "Hak Akses per Role", size=13, bold=True, color=BLUE_DEEP)
add_para(tf, "", size=3)
add_para(tf, "🔑  Admin        →  Semua fitur + manajemen user", size=10.5, color=INK, space_before=3)
add_para(tf, "👤  Karyawan  →  Aktivitas sendiri, kalender, proyek", size=10.5, color=INK, space_before=3)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 – DASHBOARD
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
section_header(slide, "04  Dashboard & Statistik",
               "Pusat informasi real-time seluruh aktivitas tim")

# Stat boxes row
stats = [
    ("124", "Total Aktivitas", BLUE_MAIN),
    ("87",  "Selesai",         GREEN),
    ("23",  "Berjalan",        AMBER),
    ("8",   "Tertunda",        RED),
    ("72%", "Rata-rata Progres", PURPLE),
]
sx = 0.3
for val, lbl, col in stats:
    stat_box(slide, val, lbl, sx, 1.38, w=2.45, h=1.2, accent=col)
    sx += 2.58

# Bottom two columns
# Left – screenshot placeholder
screenshot_placeholder(slide, 0.3, 2.78, 6.05, 4.4, "Dashboard – Chart & Statistik")

# Right – feature description
add_rect(slide, 6.55, 2.78, 6.48, 4.4, fill=WHITE, line_color=BORDER, line_width=0.75)
txb = slide.shapes.add_textbox(Inches(6.72), Inches(2.95), Inches(6.2), Inches(4.1))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
add_para(tf, "Komponen Dashboard", size=14, bold=True, color=BLUE_DEEP)
components = [
    "📊  Distribusi status aktivitas (Pie/Donut chart)",
    "📈  Tren aktivitas 7 hari terakhir (Line chart)",
    "👥  Progres per karyawan (Bar chart – Admin only)",
    "📋  Daftar aktivitas hari ini dengan progres bar",
    "📅  Reminder kalender hari ini (6 event teratas)",
    "🔔  Notifikasi & badge di navbar",
    "🔍  Filter per karyawan (khusus Admin)",
    "⚡  Real-time refresh setiap halaman di-load",
]
for c in components:
    add_para(tf, c, size=11, color=INK, space_before=5)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 – DAILY ACTIVITY (overview)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
section_header(slide, "05  Daily Activity Management",
               "Pencatatan dan pemantauan aktivitas harian karyawan")

cards = [
    ("📝", "Input Aktivitas",
     ["Judul & deskripsi aktivitas", "Tanggal pelaksanaan", "Kategori (5 pilihan)",
      "Prioritas (4 level)", "Status awal & progres", "Catatan blocker"]),
    ("📊", "Tracking Progres",
     ["Update progres 0–100%", "Auto-update status dari progres",
      "Komentar progres berseri", "Timestamp selesai otomatis",
      "Histori komentar lengkap", "Badge indikator visual"]),
    ("🔍", "Filter & Pencarian",
     ["Filter per karyawan (Admin)", "Filter per status aktivitas",
      "Pencarian teks judul", "Tampilan tabel dengan pagination",
      "Sorted by tanggal terbaru", "Export-ready layout"]),
    ("🏷️", "Kategori & Prioritas",
     ["Operasional / Administrasi", "HRGA / Koordinasi / Improvement",
      "Prioritas: Rendah → Urgent", "Status: Belum Mulai → Selesai",
      "Warna kode otomatis", "Badge per status"]),
]

cx, cy = 0.28, 1.42
for i, (icon, title, bullets) in enumerate(cards):
    col = i % 2
    row = i // 2
    lx = cx + col * 6.42
    ly = cy + row * 3.0
    feature_card(slide, icon, title, bullets, lx, ly, w=6.18, h=2.82)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 – DAILY ACTIVITY screenshots
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
section_header(slide, "05  Daily Activity – Tampilan Antarmuka",
               "List, detail, dan form input aktivitas")

screenshot_placeholder(slide, 0.3,  1.4,  6.1, 3.7, "List Aktivitas Harian")
screenshot_placeholder(slide, 6.65, 1.4,  6.38, 3.7, "Form Input Aktivitas")
screenshot_placeholder(slide, 0.3,  5.28, 6.1, 2.0, "Detail & Komentar Progres")
screenshot_placeholder(slide, 6.65, 5.28, 6.38, 2.0, "Filter & Pencarian Aktivitas")


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 – KALENDER
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
section_header(slide, "06  Kalender & Event",
               "Manajemen jadwal, meeting, dan hari libur terintegrasi")

screenshot_placeholder(slide, 0.3, 1.4, 6.5, 5.8, "Tampilan Kalender Interaktif")

add_rect(slide, 7.05, 1.4, 5.98, 5.8, fill=WHITE, line_color=BORDER, line_width=0.75)
txb = slide.shapes.add_textbox(Inches(7.2), Inches(1.55), Inches(5.7), Inches(5.5))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
add_para(tf, "Tipe Event Kalender", size=14, bold=True, color=BLUE_DEEP)
event_types = [
    ("🤝", "Meeting",   BLUE_MAIN,  "Rapat dan pertemuan tim"),
    ("🔔", "Reminder",  AMBER,      "Pengingat tugas & deadline"),
    ("🏖️", "Holiday",   GREEN,      "Hari libur nasional & perusahaan"),
    ("🏠", "Leave",     PURPLE,     "Cuti & izin karyawan"),
]
for emoji, typ, col, desc in event_types:
    add_para(tf, "", size=3)
    p = tf.add_paragraph()
    p.space_before = Pt(4)
    run = p.add_run()
    run.text = f"{emoji}  {typ}"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = col
    add_para(tf, "    " + desc, size=10.5, color=MUTED)

add_para(tf, "", size=6)
add_para(tf, "Fitur Unggulan", size=14, bold=True, color=BLUE_DEEP)
cal_features = [
    "Kalender interaktif FullCalendar.js",
    "CRUD event hanya oleh Admin",
    "Karyawan dapat melihat semua event",
    "Aktivitas harian tampil di kalender",
    "Start & end time per event",
    "Search & filter event list",
    "Warna kode otomatis per tipe",
    "Tanda khusus hari libur nasional",
]
for f in cal_features:
    add_para(tf, "✓  " + f, size=11, color=INK, space_before=4)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 – PROJECT AGILE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
section_header(slide, "07  Project Agile – Sprint & Kanban",
               "Manajemen proyek berbasis metodologi Agile/Scrum")

# Top row – 3 feature cards
cols_data = [
    ("🗂️", "Manajemen Proyek",
     ["Buat proyek dengan goal & timeline", "Status: Planning/Active/On Hold/Done",
      "Owner & member assignment", "Target date & progress tracking",
      "Visibilitas berbasis role"], BLUE_MAIN),
    ("🏃", "Sprint Management",
     ["Buat sprint dalam proyek", "Status: Planned/Active/Closed",
      "Start & end date sprint", "Sprint goal & deskripsi",
      "Multiple sprint per proyek"], PURPLE),
    ("✅", "Task & Backlog",
     ["Tipe: Story/Task/Bug/Improvement", "Prioritas: Low/Med/High/Urgent",
      "Story points estimation", "Assignee & reporter tracking",
      "Acceptance criteria"], GREEN),
]
lx = 0.28
for icon, title, bullets, accent in cols_data:
    feature_card(slide, icon, title, bullets, lx, 1.42, w=4.2, h=3.2, accent=accent)
    lx += 4.42

# Kanban workflow
add_rect(slide, 0.28, 4.8, 12.75, 2.52, fill=BLUE_SOFT, line_color=BLUE_MAIN, line_width=0.75)
add_textbox(slide, "Alur Kanban Board", 0.45, 4.88, 5, 0.38,
            size=13, bold=True, color=BLUE_DEEP)
statuses = [
    ("Backlog", MUTED),
    ("To Do",   BLUE_MAIN),
    ("In Progress", AMBER),
    ("Review",  PURPLE),
    ("Done",    GREEN),
]
kx = 0.5
for status, col in statuses:
    sw = 2.2
    add_rect(slide, kx, 5.35, sw, 1.65, fill=WHITE, line_color=col, line_width=1)
    add_rect(slide, kx, 5.35, sw, 0.38, fill=col)
    add_textbox(slide, status, kx, 5.38, sw, 0.32,
                size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Task cards\ndi-drag & drop\nper status",
                kx+0.1, 5.82, sw-0.2, 0.95,
                size=9, color=MUTED, align=PP_ALIGN.CENTER)
    if status != "Done":
        add_textbox(slide, "→", kx+sw, 5.95, 0.3, 0.4,
                    size=14, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    kx += sw + 0.36


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 – PROJECT screenshots
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
section_header(slide, "07  Project Agile – Tampilan Antarmuka",
               "Daftar proyek, detail sprint, dan kanban board")

screenshot_placeholder(slide, 0.3,  1.4,  6.1, 3.7, "Daftar Proyek")
screenshot_placeholder(slide, 6.65, 1.4,  6.38, 3.7, "Detail Proyek & Sprint")
screenshot_placeholder(slide, 0.3,  5.28, 12.73, 2.0, "Kanban Board – Status Task")


# ═══════════════════════════════════════════════════════════════
# SLIDE 12 – MANAJEMEN USER
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
section_header(slide, "08  Manajemen User (Admin)",
               "Pengelolaan akun karyawan oleh administrator")

screenshot_placeholder(slide, 0.3, 1.4, 5.6, 5.8, "Daftar & Statistik User")

add_rect(slide, 6.2, 1.4, 2.4, 1.62, fill=WHITE, line_color=BORDER, line_width=0.75)
stat_box(slide, "12", "Total Karyawan", 6.2, 1.4, w=2.4, h=1.6, accent=BLUE_MAIN)
add_rect(slide, 8.8, 1.4, 2.4, 1.62, fill=WHITE, line_color=BORDER, line_width=0.75)
stat_box(slide, "3", "Admin", 8.8, 1.4, w=2.4, h=1.6, accent=PURPLE)
add_rect(slide, 11.2, 1.4, 1.88, 1.62, fill=WHITE, line_color=BORDER, line_width=0.75)
stat_box(slide, "9", "Karyawan", 11.2, 1.4, w=1.88, h=1.6, accent=GREEN)

add_rect(slide, 6.2, 3.2, 6.88, 4.0, fill=WHITE, line_color=BORDER, line_width=0.75)
txb = slide.shapes.add_textbox(Inches(6.4), Inches(3.35), Inches(6.5), Inches(3.7))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
add_para(tf, "Fitur Manajemen User (Admin Only)", size=14, bold=True, color=BLUE_DEEP)
user_feats = [
    ("➕ Tambah User",        "Buat akun karyawan baru dengan role assignment"),
    ("✏️  Edit User",           "Ubah nama, email, dan role karyawan"),
    ("🔑 Reset Password",     "Reset password user tanpa mengetahui password lama"),
    ("📊 Statistik Aktivitas","Lihat jumlah & rata-rata progres aktivitas tiap user"),
    ("🛡️  Role Management",    "Assign role Admin atau Karyawan"),
    ("🔒 Self-Protection",    "Admin tidak bisa ubah role diri sendiri"),
]
for icon_label, desc in user_feats:
    add_para(tf, icon_label, size=12, bold=True, color=BLUE_MAIN, space_before=7)
    add_para(tf, "    " + desc, size=10.5, color=MUTED, space_before=1)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13 – RINGKASAN FITUR
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
section_header(slide, "09  Ringkasan Fitur & Manfaat",
               "Semua yang tersedia dalam satu platform terintegrasi")
add_rect(slide, 0, 1.35, 13.33, 6.15, fill=LIGHT_GRAY)

summary_cards = [
    ("📋", "Daily Activity",      "Catat, lacak, dan update progres aktivitas harian dengan mudah",           BLUE_MAIN),
    ("📅", "Kalender",            "Event meeting, reminder, cuti, dan hari libur dalam satu tampilan",         GREEN),
    ("🚀", "Project Agile",       "Manajemen proyek dengan sprint dan kanban board ala metodologi Scrum",     PURPLE),
    ("📊", "Dashboard",           "Statistik real-time, grafik performa, dan ringkasan tim untuk monitoring", AMBER),
    ("👥", "Manajemen User",      "Admin kelola seluruh akun karyawan, reset password, dan pantau performa", RED),
    ("🔔", "Notifikasi",          "Badge dan notifikasi navbar untuk pengingat event dan update aktivitas",   RGBColor(0x06, 0x96, 0x88)),
]

cx, cy = 0.28, 1.55
for i, (icon, title, desc, accent) in enumerate(summary_cards):
    col = i % 3
    row = i // 3
    lx = cx + col * 4.37
    ly = cy + row * 2.85
    add_rect(slide, lx, ly, 4.22, 2.62, fill=WHITE, line_color=BORDER, line_width=0.75)
    add_rect(slide, lx, ly, 4.22, 0.06, fill=accent)
    add_textbox(slide, icon + "  " + title, lx+0.15, ly+0.18, 4.0, 0.45,
                size=14, bold=True, color=accent)
    add_textbox(slide, desc, lx+0.15, ly+0.72, 3.92, 1.65,
                size=10.5, color=INK, wrap=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 14 – PENUTUP
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
gradient_bg(slide)
add_rect(slide, 0, 0, 13.33, 7.5, fill=None)
add_rect(slide, 0, 5.5, 13.33, 2.0, fill=RGBColor(0x05, 0x28, 0x4d))

# Decorative
add_rect(slide, 9.8, -0.8, 4.5, 4.5, fill=RGBColor(0x0a, 0x48, 0x8c))
add_rect(slide, -0.8, 4.2, 3.5, 3.5, fill=RGBColor(0x06, 0x30, 0x5e))

add_textbox(slide, "Terima Kasih", 1.5, 1.5, 10.33, 1.2,
            size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "HRGA Activity – Solusi Digital untuk Produktivitas Tim",
            1.5, 2.85, 10.33, 0.55,
            size=16, color=RGBColor(0xc7, 0xdc, 0xf8), align=PP_ALIGN.CENTER)
add_rect(slide, 4, 3.6, 5.33, 0.04, fill=RGBColor(0x4d, 0x96, 0xe0))
add_textbox(slide, "PT Sulawesi Cahaya Mineral  •  Departemen HRGA",
            1.5, 3.78, 10.33, 0.45,
            size=13, color=RGBColor(0xb8, 0xd4, 0xf8), align=PP_ALIGN.CENTER)

add_textbox(slide, "📧  eddyyucca@gmail.com   |   🌐  HRGA Activity v1.0",
            1.5, 5.62, 10.33, 0.45,
            size=11, color=RGBColor(0x8b, 0xb4, 0xd8), align=PP_ALIGN.CENTER)
add_textbox(slide, "© 2024 PT Sulawesi Cahaya Mineral. All rights reserved.",
            1.5, 6.18, 10.33, 0.4,
            size=10, color=RGBColor(0x64, 0x7a, 0x96), align=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────
out = r"c:\xampp\htdocs\timeline_kalender_hrga\HRGA_Activity_Presentation.pptx"
prs.save(out)
print("Saved:", out)
